import pandas as pd
from PIL import Image
import streamlit as st
import json
import datetime
import os
import glob
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

combined_df = []
show_overall = False

def print_text_box(pSlide, pTopNo, pVal, pBold, pColor, pSize):
    left = Inches(0.5)
    top = Inches(pTopNo)
    width = Inches(6.0)
    height = Inches(0.3)
    text_box = pSlide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = pVal
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(pSize)
    if (pBold):
        font.bold = True
    if (pColor):
        font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

def clear_button():
    try:
        # Add exception handling for potential errors within this function
        edited_df = pd.DataFrame(empty_df)
        session_state.clear()
    except Exception as e:
        st.error(f"An error occurred while clearing the button: {e}")

def show_user_report():
    try:
        # Add exception handling for potential errors within this function
        file_name = f"Weekly_Files/{selected_location}_{selected_name}_{selected_week}.csv"
        session_state.df = pd.read_csv(file_name)
    except Exception as e:
        st.error(f"An error occurred while showing the user report: {e}")

def generate_ppt():
    try:
        # Add exception handling for potential errors within this function
        df = pd.DataFrame(combined_df)
        grouped = df.groupby('Category')
        descriptions = df.groupby('Description')
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Weekly Report CW{selected_week}"  
        print_text_box(slide, 1.2, f"{presentername} \n {current_date}", True, True, 18)
        for sildename, group in grouped:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5]) 
            slide.shapes.title.text = f"Weekly Meeting CW{selected_week}"  
            leftNo = 0.5
            topNo = 1.2
            print_text_box(slide, topNo, sildename, True, True, 20)
            topics = group.groupby('Topic')
            heightNo = 0.3
            for name, topic in topics:
                print(name)
                topNo += heightNo
                print_text_box(slide, topNo, name, True, False, 18)
                heightNo = 0.3
                for d in topic['Description']:
                    heightNo += 0.3
                    left = Inches(0.5)
                    top = Inches(topNo + 0.3)
                    width = Inches(6.0)
                    height = Inches(heightNo)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    p = text_frame.paragraphs[0]
                    print(d)
                    topNo += 0.3
                    run = p.add_run()
                    run.text = "- " + d
                    p.level = 1
                    p.add_line_break()
                    font = run.font
                    font.name = 'Calibri'
                    font.size = Pt(12)
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  
        presentation.save(f'{pptfilename}{current_date}.pptx')

    except Exception as e:
        st.error(f"An error occurred while showing the user report: {e}")

if __name__ == "__main__":
    try:
        st.set_page_config(
            page_title="Weekly Meeting",
            page_icon="👋 Hello EA",
            layout='wide'
        )
        st.write("**Welcome to Weekly Reports Automation Tool! 👋**")
        calendar_weeks = [int(datetime.datetime.utcnow().isocalendar()[1]) -1, datetime.datetime.utcnow().isocalendar()[1], int(datetime.datetime.utcnow().isocalendar()[1]) +1]
        current_date = datetime.datetime.utcnow().date()
        try:
            with open('config.json', 'r') as file:
                config = json.load(file)
                locations = config['locations']
                names = config['names']
                category = config['category']
                topics = config['topic']
                pptfilename = config['outputfilename']
                presentername = config['presentername']
        except Exception as e:
            st.error(f"An error occurred in the while reading config file: {e}")

        col1, col2, col3, col4, col5 = st.columns([1, 1, 1, 1, 1])

        with col1:
            selected_location = st.sidebar.selectbox('Select Location', locations, on_change=clear_button)

        with col2:
            selected_name = st.sidebar.selectbox('Select Name', names[selected_location], on_change=clear_button)

        with col3:
            selected_week = st.sidebar.selectbox('Select Calendar Week', calendar_weeks, on_change=clear_button)

        with col4:
            col1, col2, col3 = st.columns(3)
            with col1:
                if st.sidebar.button("Show Overall", type="primary"):
                    file_pattern = f'Weekly_Files/{selected_location}*{selected_week}.csv'  
                    csv_files = glob.glob(file_pattern)
                    dfs = []
                    for file in csv_files:
                        try:
                            df = pd.read_csv(file)
                        except Exception as e:
                            st.error(f"An error occurred in the while reading individual file for combined report: {e}")
                        dfs.append(df)
                    combined_df = pd.concat(dfs, ignore_index=True)
                    show_overall = True
            with col2:
                if st.sidebar.button("Hide Overall", type="primary"):
                    show_overall = False
            with col3:
                if st.sidebar.button("Generte PPT", type="primary", on_click=generate_ppt):
                    show_overall = False
                
        st.write('You selected:', selected_location, 'for location,', selected_name, 'for reports, and', selected_week, 'th for calendar week.')

        st.write('                                                    ')
        st.write('                                                    ')
        selected_category = st.selectbox('Select Category', category)
        selected_topic = st.selectbox('Select Topic', topics[selected_location])
        Description = st.text_input('Enter Description', 'Please enter your update here')
        columns = ["Category","Topic","Description", "Name"]
        empty_df = pd.DataFrame(columns=columns)
        session_state = st.session_state

        if "df" not in session_state:
            session_state.df = empty_df
            session_state.row = pd.Series(index=columns)
        if st.button("Add Update"):
            new_data = {'Category': selected_category, 'Topic': selected_topic, 'Description' : Description, 'Name' : selected_name}
            session_state.df = session_state.df._append(new_data, ignore_index=True)
            session_state.row = pd.Series(index=columns)
            st.success("Entry updated successfully.")


        edited_df = st.data_editor(session_state.df, use_container_width= True)
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("Save it !", type= "primary"):
                file_name = f"Weekly_Files/{selected_location}_{selected_name}_{selected_week}.csv"
                edited_df.to_csv(file_name, index=False)
        with col2:
            if st.button("Clear update", type= "primary", on_click=clear_button):
                pass
        with col3:
            if st.button("Show Mine", type= "primary", on_click=show_user_report):
                pass

        st.write('                                                    ')
        st.write('                                                    ')
        st.write('                                                    ')
        st.write('                                                    ')
        st.write('                                                    ')
        st.write('                                                    ')
        st.write('                                                    ')
        st.write('                                                    ')

        if show_overall == True:
            st.write("*Overall Reports*")
            st.data_editor(combined_df, use_container_width=True, key = "Overall")
    except Exception as e:
        st.error(f"An error occurred in the main block: {e}")



